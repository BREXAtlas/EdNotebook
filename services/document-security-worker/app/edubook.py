from __future__ import annotations

import re
import uuid
from pathlib import Path
from typing import Any, Iterable

import ebooklib
from bs4 import BeautifulSoup
from docx import Document
from ebooklib import epub
from pptx import Presentation
from pypdf import PdfReader


MAX_MANIFEST_CHARACTERS = 1_500_000
MAX_BLOCK_CHARACTERS = 12_000


def _id() -> str:
    return str(uuid.uuid4())


def _clean(value: str) -> str:
    return re.sub(r"\s+", " ", value or "").strip()


def _paragraph_blocks(text: str) -> list[dict[str, Any]]:
    blocks: list[dict[str, Any]] = []
    for paragraph in re.split(r"\n\s*\n|(?<=\.)\s*\n", text):
        cleaned = _clean(paragraph)
        if not cleaned:
            continue
        for offset in range(0, len(cleaned), MAX_BLOCK_CHARACTERS):
            blocks.append({"id": _id(), "type": "paragraph", "text": cleaned[offset:offset + MAX_BLOCK_CHARACTERS]})
    return blocks


def _chapter(title: str, blocks: list[dict[str, Any]], locator: str | None = None) -> dict[str, Any]:
    return {
        "id": _id(),
        "title": _clean(title) or "Untitled chapter",
        "locator": locator,
        "blocks": blocks,
        "knowledgeChecks": [],
        "discussionPrompts": [],
    }


def _text_chapters(text: str) -> list[dict[str, Any]]:
    lines = text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    chapters: list[dict[str, Any]] = []
    current_title = "Opening"
    current_lines: list[str] = []

    def commit() -> None:
        nonlocal current_lines
        blocks = _paragraph_blocks("\n".join(current_lines))
        if blocks or not chapters:
            chapters.append(_chapter(current_title, blocks))
        current_lines = []

    for line in lines:
        stripped = line.strip()
        heading = re.match(r"^#{1,4}\s+(.+)$", stripped)
        numbered = re.match(r"^(?:chapter|unit|part|section)\s+[0-9ivxlcdm]+[\s:.-]+(.+)$", stripped, re.I)
        if heading or numbered:
            if current_lines:
                commit()
            current_title = (heading or numbered).group(1).strip()
        else:
            current_lines.append(line)
    commit()
    return chapters


def _docx_chapters(path: Path) -> list[dict[str, Any]]:
    document = Document(str(path))
    chapters: list[dict[str, Any]] = []
    title = "Opening"
    blocks: list[dict[str, Any]] = []

    def commit() -> None:
        nonlocal blocks
        if blocks or not chapters:
            chapters.append(_chapter(title, blocks))
        blocks = []

    for paragraph in document.paragraphs:
        text = _clean(paragraph.text)
        if not text:
            continue
        style = (paragraph.style.name if paragraph.style else "").lower()
        if style.startswith("heading") or style in {"title", "subtitle"}:
            if blocks:
                commit()
            title = text
        else:
            blocks.append({"id": _id(), "type": "paragraph", "text": text[:MAX_BLOCK_CHARACTERS]})

    for table_index, table in enumerate(document.tables, start=1):
        rows = [[_clean(cell.text) for cell in row.cells] for row in table.rows]
        blocks.append({
            "id": _id(),
            "type": "table",
            "title": f"Table {table_index}",
            "rows": rows[:200],
        })
    commit()
    return chapters


def _pptx_chapters(path: Path) -> list[dict[str, Any]]:
    presentation = Presentation(str(path))
    chapters: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        title = ""
        if slide.shapes.title:
            title = _clean(slide.shapes.title.text)
        text_items: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text") or shape is slide.shapes.title:
                continue
            text = _clean(shape.text)
            if text:
                text_items.append(text)
        notes = ""
        try:
            notes = _clean(slide.notes_slide.notes_text_frame.text)
        except Exception:
            notes = ""
        blocks = [{"id": _id(), "type": "paragraph", "text": value[:MAX_BLOCK_CHARACTERS]} for value in text_items]
        if notes:
            blocks.append({"id": _id(), "type": "speaker-notes", "text": notes[:MAX_BLOCK_CHARACTERS]})
        chapters.append(_chapter(title or f"Slide {index}", blocks, locator=f"slide:{index}"))
    return chapters


def _pdf_chapters(path: Path) -> list[dict[str, Any]]:
    reader = PdfReader(str(path), strict=False)
    chapters: list[dict[str, Any]] = []
    total_characters = 0
    for index, page in enumerate(reader.pages, start=1):
        if total_characters >= MAX_MANIFEST_CHARACTERS:
            break
        text = (page.extract_text() or "").strip()
        if not text:
            chapters.append(_chapter(f"Page {index}", [], locator=f"page:{index}"))
            continue
        allowed = text[: MAX_MANIFEST_CHARACTERS - total_characters]
        total_characters += len(allowed)
        first_line = next((line.strip() for line in allowed.splitlines() if line.strip()), "")
        title = first_line[:120] if 3 <= len(first_line) <= 120 else f"Page {index}"
        chapters.append(_chapter(title, _paragraph_blocks(allowed), locator=f"page:{index}"))
    return chapters


def _epub_chapters(path: Path) -> list[dict[str, Any]]:
    book = epub.read_epub(str(path), options={"ignore_ncx": False})
    chapters: list[dict[str, Any]] = []
    total_characters = 0
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        if total_characters >= MAX_MANIFEST_CHARACTERS:
            break
        soup = BeautifulSoup(item.get_content(), "html.parser")
        heading = soup.find(["h1", "h2", "h3", "title"])
        title = _clean(heading.get_text(" ", strip=True)) if heading else Path(item.get_name()).stem
        text = soup.get_text("\n\n", strip=True)
        allowed = text[: MAX_MANIFEST_CHARACTERS - total_characters]
        total_characters += len(allowed)
        chapters.append(_chapter(title or f"Chapter {len(chapters) + 1}", _paragraph_blocks(allowed), locator=item.get_name()))
    return chapters


def _truncate_manifest(chapters: Iterable[dict[str, Any]]) -> list[dict[str, Any]]:
    result: list[dict[str, Any]] = []
    consumed = 0
    for chapter in chapters:
        kept_blocks: list[dict[str, Any]] = []
        for block in chapter.get("blocks", []):
            serialized_size = len(str(block))
            if consumed + serialized_size > MAX_MANIFEST_CHARACTERS:
                break
            consumed += serialized_size
            kept_blocks.append(block)
        result.append({**chapter, "blocks": kept_blocks})
        if consumed >= MAX_MANIFEST_CHARACTERS:
            break
    return result


def build_edubook(
    path: Path,
    mime_type: str,
    *,
    title: str,
    author: str = "",
    description: str = "",
    checksum_sha256: str | None = None,
    original_name: str | None = None,
) -> dict[str, Any]:
    if mime_type in {"text/plain", "text/markdown", "text/html"}:
        source_text = path.read_text("utf-8", errors="replace")[:MAX_MANIFEST_CHARACTERS]
        if mime_type == "text/html":
            source_text = BeautifulSoup(source_text, "html.parser").get_text("\n\n", strip=True)
        chapters = _text_chapters(source_text)
    elif mime_type in {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    }:
        chapters = _docx_chapters(path)
    elif mime_type in {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }:
        chapters = _pptx_chapters(path)
    elif mime_type == "application/pdf":
        chapters = _pdf_chapters(path)
    elif mime_type == "application/epub+zip":
        chapters = _epub_chapters(path)
    else:
        raise ValueError(f"EduBook conversion is not supported for {mime_type}")

    chapters = _truncate_manifest(chapters)
    word_count = sum(
        len(str(block.get("text", "")).split())
        for chapter in chapters
        for block in chapter.get("blocks", [])
    )
    return {
        "format": "EduBook/1.0",
        "title": _clean(title) or Path(original_name or path.name).stem,
        "author": _clean(author) or "Unknown author",
        "description": _clean(description),
        "language": "en",
        "source": {
            "type": mime_type,
            "originalName": original_name or path.name,
            "checksumSha256": checksum_sha256,
            "convertedAt": __import__("datetime").datetime.now(__import__("datetime").timezone.utc).isoformat(),
            "words": word_count,
        },
        "learningDesign": {
            "mode": "interactive-reading",
            "subjectAgnostic": True,
            "annotations": True,
            "bookmarks": True,
            "progress": True,
            "checks": True,
            "discussion": True,
            "teacherLayerSeparatedFromSource": True,
        },
        "chapters": chapters,
    }
